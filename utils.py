import streamlit as st
import torch
from transformers import BartTokenizer, BartForConditionalGeneration
from docx import Document
from pptx import Presentation
import PyPDF2


# ==============================
# LOAD MODEL (CACHED)
# ==============================

@st.cache_resource
def load_model():
    model_name = "sshleifer/distilbart-cnn-12-6"

    tokenizer = BartTokenizer.from_pretrained(model_name)
    model = BartForConditionalGeneration.from_pretrained(model_name)

    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    model.to(device)

    return tokenizer, model, device


# ==============================
# FILE READERS
# ==============================

def read_pdf(file):
    text = ""
    reader = PyPDF2.PdfReader(file)
    for page in reader.pages:
        text += page.extract_text() + " "
    return text


def read_docx(file):
    text = ""
    doc = Document(file)
    for p in doc.paragraphs:
        text += p.text + " "
    return text


def read_ppt(file):
    text = ""
    prs = Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text


def read_txt(file):
    return file.read().decode("utf-8")


# ==============================
# SUMMARY GENERATION
# ==============================

def generate_summary(text, length_option="Medium"):
    tokenizer, model, device = load_model()

    # Define lengths based on option
    if length_option == "Short":
        max_l, min_l = 150, 50
    elif length_option == "Long":
        max_l, min_l = 500, 200
    else:  # Medium
        max_l, min_l = 300, 120

    inputs = tokenizer(
        text,
        max_length=1024,
        return_tensors="pt",
        truncation=True
    ).to(device)

    summary_ids = model.generate(
        inputs["input_ids"],
        max_length=max_l,
        min_length=min_l,
        num_beams=4,
        no_repeat_ngram_size=3,
        early_stopping=True
    )

    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return format_summary(summary)


# ==============================
# FORMAT SUMMARY
# ==============================

def format_summary(text):
    sentences = text.split(". ")
    formatted = ""

    for s in sentences:
        s = s.strip()
        if s:
            if not s.endswith("."):
                s += "."
            formatted += "• " + s + "\n\n"

    return formatted